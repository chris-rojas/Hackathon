import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os


# Sample data for financials and geographic distribution
data = {
    "Region": ["North America", "Europe", "Asia"],
    "Revenue": [120000, 85000, 100000],
    "Profit": [20000, 15000, 18000]
}

df = pd.DataFrame(data)

def create_bar_chart():
    fig, ax = plt.subplots()
    ax.bar(df["Region"], df["Revenue"], color="skyblue")
    ax.set_title("Revenue by Region")
    ax.set_xlabel("Region")
    ax.set_ylabel("Revenue ($)")
    fig_path = "result/bar_chart.png"
    plt.savefig(fig_path)  # Save the chart as an image file
    plt.close(fig)
    return fig_path

# Function to save content to a PowerPoint file
def save_content_to_ppt(content, chart_path, filename="result/Profiler_Report.pptx"):
    # Ensure the result directory exists
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    
    # Create a PowerPoint presentation
    prs = Presentation()
    
    for title, text in content.items():
        # Add a slide for each section
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        # Set the title and content
        title_placeholder.text = title
        content_placeholder.text = text

        # Add the chart to the Financials slide
        if title == "Financials" and os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(1), Inches(2), width=Inches(4))

    # Save the presentation
    prs.save(filename)


def show_profiler():
    # Load custom CSS
    with open("style.css") as f:
        st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

    # Profiler sub-tabs
    profiler_tabs = st.tabs(["Overview", "Financials", "Geographic Mix", "Management", "Recent News"])

    # Collect content for each tab
    content = {}
    with profiler_tabs[0]:
        content["Overview"] = "This is the company overview with a sample paragraph.\nLorem ipsum dolor sit amet, consectetur adipiscing elit. Praesent nec nisl vel mauris blandit interdum."
        st.write("### Overview")
        st.write(content["Overview"])

    with profiler_tabs[1]:
        content["Financials"] = "Financial Table:\n" + df.to_string()
        st.write("### Financial Table")
        st.dataframe(df)
        st.write("### Financial Bar Chart")
        chart_path = create_bar_chart()
        st.image(chart_path)

    with profiler_tabs[2]:
        content["Geographic Mix"] = "Geographic Distribution Table:\n" + df[["Region", "Revenue"]].to_string()
        st.write("### Geographic Distribution Table")
        st.dataframe(df[["Region", "Revenue"]])

    with profiler_tabs[3]:
        content["Management"] = "This section could include profiles or key information about the management team."
        st.write("### Management Information")
        st.write(content["Management"])

    with profiler_tabs[4]:
        content["Recent News"] = "1. Company launches new product.\n2. Q3 earnings report shows positive growth.\n3. Expansion into new markets."
        st.write("### Recent News")
        st.write(content["Recent News"])



    # Button to generate and download the PowerPoint file
    if st.button("Download PPT"):
        save_content_to_ppt(content, chart_path)  # Save the content to a PowerPoint file
#         with open("result/Profiler_Report.pptx", "rb") as file:
#             btn = st.download_button(
#                 label="Download PowerPoint",
#                 data=file,
#                 file_name="Profiler_Report.pptx",
#                 mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#             )
